import subprocess
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import argparse
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# --- Output folder ---
PLOTS_FOLDER = "plots"
os.makedirs(PLOTS_FOLDER, exist_ok=True)

def pretty_table(df, title, filename, col_width=2.8, row_height=0.8):
    import matplotlib.pyplot as plt
    import numpy as np

    n_rows, n_cols = df.shape
    fig_width = col_width * n_cols
    fig_height = row_height * (n_rows + 1)  # +1 for header row

    plt.figure(figsize=(fig_width, fig_height))
    plt.axis('off')

    # Create the table
    tbl = plt.table(
        cellText=df.values,
        colLabels=df.columns,
        loc='center',
        cellLoc='left',
        colLoc='left',
        edges='open',
        bbox=[0, 0, 1, 1]
    )

    tbl.auto_set_font_size(False)
    tbl.set_fontsize(14)

    # Style header row (row 0)
    for j in range(n_cols):
        cell = tbl[(0, j)]
        cell.set_text_props(weight='bold', color='white')
        cell.set_facecolor('#2E4053')  # deep blue for header
        cell.PAD = 0.12

    # Adjust first column cells: increase padding to reduce overlap
    first_col_width_adjustment = 0.18

    # Style data rows and specifically adjust first column
    for i in range(1, n_rows + 1):
        for j in range(n_cols):
            cell = tbl[(i, j)]
            # Alternating row color
            cell.set_facecolor('#F2F4F4' if i % 2 == 0 else 'white')
            cell.set_text_props(color='black', weight='normal')
            # Increase padding for first column to reduce overlap
            if j == 0:
                cell.PAD = 0.18
            else:
                cell.PAD = 0.12

    plt.title(title, fontsize=20, pad=16, loc='left')

    # Try to auto adjust columns (especially widen first column)
    tbl.auto_set_column_width(col=list(range(n_cols)))

    outpath = os.path.join(PLOTS_FOLDER, filename)
    plt.savefig(outpath, bbox_inches='tight', dpi=200)
    plt.close()
    print(f"Saved pretty table image: {outpath}")
    return filename


def plot_grouped_bar(df, metric, title, ylabel, filename, palette="tab20"):
    plt.figure(figsize=(16, 9))
    sns.barplot(data=df, x="streams", y=metric, hue="method",
                palette=palette, edgecolor="black")
    plt.title(title, fontsize=20, loc='left')
    plt.xlabel("Streams", fontsize=14)
    plt.ylabel(ylabel, fontsize=14)
    plt.legend(title="Method", loc="best", fontsize=12)
    plt.tight_layout()
    save_path = os.path.join(PLOTS_FOLDER, filename)
    plt.savefig(save_path)
    plt.close()
    print(f"Saved grouped bar chart: {save_path}")

def plot_metric(df, metric, title, ylabel, filename, palette="viridis"):
    plt.figure(figsize=(16, 9))
    sns.barplot(data=df, x="method", y=metric, hue="method", palette=palette, legend=False)
    plt.title(title, fontsize=20, loc='left')
    plt.xlabel("Method", fontsize=14)
    plt.ylabel(ylabel, fontsize=14)
    plt.xticks(rotation=30, ha='right', fontsize=12)
    plt.yticks(fontsize=12)
    plt.tight_layout()
    save_path = os.path.join(PLOTS_FOLDER, filename)
    plt.savefig(save_path)
    plt.close()
    print(f"Saved plot: {save_path}")

def plot_scaling(df, metric, title, ylabel, filename, legend_loc="best"):
    plt.figure(figsize=(16, 9))
    sns.lineplot(data=df, x="streams", y=metric, hue="method", marker="o")
    plt.title(title, fontsize=20, loc='left')
    plt.xlabel("Streams", fontsize=14)
    plt.ylabel(ylabel, fontsize=14)
    plt.legend(title="Method", loc=legend_loc, fontsize=12)
    plt.tight_layout()
    save_path = os.path.join(PLOTS_FOLDER, filename)
    plt.savefig(save_path)
    plt.close()
    print(f"Saved plot: {save_path}")

def blank_image(filename):
    plt.figure(figsize=(10, 2))
    plt.axis("off")
    outpath = os.path.join(PLOTS_FOLDER, filename)
    plt.savefig(outpath, bbox_inches='tight', dpi=100)
    plt.close()
    return filename

def save_to_ppt(slides, ppt_filename):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only layout

        # Remove default title placeholder to avoid overlap
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 1:
                slide.shapes._spTree.remove(shape._element)

        # Add custom full-width title text box, left aligned
        title_box = slide.shapes.add_textbox(
            0, Inches(0.15), prs.slide_width, Inches(1)
        )
        title_tf = title_box.text_frame
        title_tf.clear()
        title_p = title_tf.paragraphs[0]
        title_p.text = slide_data['title']
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.alignment = PP_ALIGN.LEFT

        if slide_data.get('subtitle'):
            left = Inches(0.5)
            top = Inches(1.1)
            width = prs.slide_width - Inches(1)
            height = Inches(0.6)
            subtitle_box = slide.shapes.add_textbox(left, top, width, height)
            tf = subtitle_box.text_frame
            tf.text = slide_data['subtitle']
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        img_path = os.path.join(PLOTS_FOLDER, slide_data['filename'])
        if os.path.isfile(img_path):
            left_img = Inches(0.5)
            top_img = Inches(1.9)
            max_width = prs.slide_width - Inches(1)
            max_height = prs.slide_height - Inches(2.1)
            slide.shapes.add_picture(img_path, left_img, top_img, width=max_width, height=max_height)
    ppt_out = os.path.join(PLOTS_FOLDER, ppt_filename)
    prs.save(ppt_out)
    print(f"\n✅ PowerPoint file created: {ppt_out}")

def generate_stream_runs(max_streams):
    base = [x for x in [1, 3, 5] if x <= max_streams]
    if max_streams > 5:
        base += list(range(10, max_streams + 1, 5))
    return base

def run_benchmark(input_file, streams, exe='./benchmark_all_9'):
    print(f"\n▶️  Running benchmark with {streams} streams...")
    result = subprocess.run(
        [exe, input_file, str(streams)],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        encoding='utf-8'
    )
    if result.returncode != 0:
        print(f"Error running benchmark: {result.stderr}")
        return pd.DataFrame(), result.stdout
    return parse_output(result.stdout, streams), result.stdout

def parse_output(output_text, stream_count):
    in_table = False
    results = []
    for line in output_text.split('\n'):
        line = line.strip()
        if 'Method' in line and 'Time/Frame' in line:
            in_table = True
            continue
        if in_table:
            if line.startswith("—") or line == "" or line.startswith("---"):
                continue
            parts = [x.strip() for x in line.split('|')]
            if len(parts) < 8:
                continue
            try:
                method = parts[0]
                time_per_frame = float(parts[1].replace("ms", "").strip())
                fps = float(parts[2])
                cpu = float(parts[3].replace("%", ""))
                mem = float(parts[4])
                mvs = int(parts[5])
                frames = int(parts[6])
                high_profile = parts[7]
                results.append({
                    "method": method,
                    "streams": stream_count,
                    "time_per_frame": time_per_frame,
                    "fps": fps,
                    "cpu": cpu,
                    "memory": mem,
                    "mvs": mvs,
                    "frames": frames,
                    "high_profile": high_profile
                })
            except Exception:
                pass
    return pd.DataFrame(results)

def run_all(input_path, max_streams, exe_path):
    stream_steps = generate_stream_runs(max_streams)
    print(f"🔁 Stream ranges to test: {stream_steps}")

    all_results = []
    for s in stream_steps:
        df, _ = run_benchmark(input_path, s, exe=exe_path)
        if df.empty:
            print(f"Warning: No data returned for streams={s}")
        all_results.append(df)

    full_df = pd.concat(all_results, ignore_index=True)

    exclude_methods = ["LIVE555 Parser", "Custom H.264 Parser"]
    full_df = full_df[~full_df['method'].isin(exclude_methods)].copy()

    csv_path = os.path.join(PLOTS_FOLDER, "benchmark_results.csv")
    full_df.to_csv(csv_path, index=False)
    print(f"Saved complete data table: {csv_path}")

    df_hp = full_df[full_df["high_profile"] == "✅"].copy()
    if df_hp.empty:
        print("No high profile algorithms found in results!")
        return full_df

    slides = []

    # 1. Fastest Methods Table
    streams_order = sorted(df_hp["streams"].unique())
    rows = []
    for s in streams_order:
        sub = df_hp[(df_hp["streams"] == s)]
        if not sub.empty:
            fastest = sub.loc[sub["time_per_frame"].idxmin()]
            rows.append([
                s, fastest["method"], fastest["time_per_frame"], fastest["fps"], fastest["cpu"]
            ])
    tbl_fastest = pd.DataFrame(rows, columns=["Streams", "Method", "Time/Frame (ms)", "FPS", "CPU (%)"])
    pretty_table(tbl_fastest, "Fastest High Profile Method Per Streams", "fastest_high_profile_methods.png")
    slides.append({
        "title": "Fastest Methods",
        "subtitle": "Best (lowest time/frame) method at each streams value",
        "filename": "fastest_high_profile_methods.png"
    })

    # 2. Scaling (Line) Charts
    scaling_plots = [
        dict(metric="fps", title="Throughput Scaling", ylabel="Frames per Second (Higher = Better)",
             filename="scaling_fps.png", subtitle="High Profile Methods: FPS vs Streams"),
        dict(metric="time_per_frame", title="Latency Scaling", ylabel="Time per Frame (ms, Lower = Better)",
             filename="scaling_timeperframe.png", subtitle="High Profile Methods: Time per Frame vs Streams"),
        dict(metric="cpu", title="CPU Usage Scaling", ylabel="CPU Usage (%)",
             filename="scaling_cpu.png", subtitle="High Profile Methods: CPU Usage (%) vs Streams"),
        dict(metric="memory", title="Memory Usage Scaling", ylabel="Memory (kB)",
             filename="scaling_memory.png", subtitle="High Profile Methods: Memory Usage (kB) vs Streams")
    ]
    for cfg in scaling_plots:
        plot_scaling(df_hp, cfg["metric"], cfg["title"] + "", cfg["ylabel"], cfg["filename"])
        slides.append({
            "title": cfg["title"] + "",
            "subtitle": cfg["subtitle"],
            "filename": cfg["filename"]
        })

    # 3. Grouped Bar Charts
    plot_grouped_bar(df_hp, "fps", "Algorithm Throughput (FPS) vs Streams — Grouped Bar",
                     "Frames per Second (Higher = Better)", "grouped_barchart_fps.png")
    slides.append({
        "title": "Grouped FPS Comparison (All Streams)",
        "subtitle": "All High Profile Methods: FPS per Streams, Grouped Bar Chart",
        "filename": "grouped_barchart_fps.png"
    })
    plot_grouped_bar(df_hp, "time_per_frame",
                     "Algorithm Latency (ms/frame) vs Streams — Grouped Bar",
                     "Time per Frame (ms, Lower = Better)", "grouped_barchart_timeperframe.png")
    slides.append({
        "title": "Grouped Latency Comparison (All Streams)",
        "subtitle": "All High Profile Methods: Latency (ms/frame) per Streams, Grouped Bar Chart",
        "filename": "grouped_barchart_timeperframe.png"
    })
    plot_grouped_bar(df_hp, "cpu", "Algorithm CPU Usage (%) vs Streams — Grouped Bar",
                     "CPU Usage (%)", "grouped_barchart_cpu.png")
    slides.append({
        "title": "Grouped CPU Usage Comparison (All Streams)",
        "subtitle": "All High Profile Methods: CPU Usage per Streams, Grouped Bar Chart",
        "filename": "grouped_barchart_cpu.png"
    })
    plot_grouped_bar(df_hp, "memory", "Algorithm Memory Usage (kB) vs Streams — Grouped Bar",
                     "Memory Usage (kB)", "grouped_barchart_memory.png")
    slides.append({
        "title": "Grouped Memory Usage Comparison (All Streams)",
        "subtitle": "All High Profile Methods: Memory Usage per Streams, Grouped Bar Chart",
        "filename": "grouped_barchart_memory.png"
    })

    # 4. Section Header for Detailed Tables
    blank_img = blank_image("blank.png")
    slides.append({
        "title": "Detailed Tables",
        "subtitle": "Full Per-Streams Benchmark Results",
        "filename": "blank.png"
    })

    # 5. Detailed Tables per Streams Count
    for streams in streams_order:
        df_sub = df_hp[df_hp["streams"] == streams]
        tbl = df_sub[["method", "time_per_frame", "fps", "cpu", "memory", "mvs", "frames"]].copy()
        tbl.columns = ["Method", "Time/frame (ms)", "FPS", "CPU (%)", "Mem Δ KB", "Total MVs", "Frames"]
        tbl_filename = f"detail_table_{streams}streams.png"
        pretty_table(tbl, f"All Methods: Streams={streams}", tbl_filename)
        slides.append({
            "title": f"Detailed Table: Streams={streams}",
            "subtitle": f"All metrics for high-profile methods, streams={streams}",
            "filename": tbl_filename
        })

    # 6. Individual Bar Charts per Streams & Metric
    for streams in streams_order:
        df_sub = df_hp[df_hp['streams'] == streams]
        fname_fps = f"barchart_fps_{streams}streams.png"
        plot_metric(df_sub, "fps", f"Algorithm Comparison: FPS @ {streams} Streams",
                    "Frames per Second (Higher = Better)", fname_fps, "viridis")
        slides.append({
            "title": f"Algorithm FPS Comparison ({streams} Streams)",
            "subtitle": f"Throughput (FPS) by High Profile Method at {streams} Streams",
            "filename": fname_fps
        })
        fname_latency = f"barchart_timeperframe_{streams}streams.png"
        plot_metric(df_sub, "time_per_frame", f"Algorithm Comparison: Time/Frame @ {streams} Streams",
                    "Time per Frame (ms, Lower = Better)", fname_latency, "mako")
        slides.append({
            "title": f"Algorithm Latency Comparison ({streams} Streams)",
            "subtitle": f"Time per Frame (ms) by High Profile Method at {streams} Streams",
            "filename": fname_latency
        })
        fname_cpu = f"barchart_cpu_{streams}streams.png"
        plot_metric(df_sub, "cpu", f"Algorithm Comparison: CPU % @ {streams} Streams",
                    "CPU Usage (%)", fname_cpu, "rocket")
        slides.append({
            "title": f"Algorithm CPU Usage Comparison ({streams} Streams)",
            "subtitle": f"CPU Usage (%) by High Profile Method at {streams} Streams",
            "filename": fname_cpu
        })
        fname_mem = f"barchart_memory_{streams}streams.png"
        plot_metric(df_sub, "memory", f"Algorithm Comparison: Memory @ {streams} Streams",
                    "Memory (kB)", fname_mem, "crest")
        slides.append({
            "title": f"Algorithm Memory Usage Comparison ({streams} Streams)",
            "subtitle": f"Memory Peak (kB) by High Profile Method at {streams} Streams",
            "filename": fname_mem
        })

    save_to_ppt(slides, "benchmark_comparison_slides_high_profile.pptx")
    return full_df

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Benchmark, save all charts, and auto-create PowerPoint slides (High Profile only).")
    parser.add_argument("input", help="Input video file or RTSP URL")
    parser.add_argument("streams", type=int, help="Maximum stream count")
    parser.add_argument("--exe", default="./benchmark_all_9", help="Benchmark executable to run")
    args = parser.parse_args()
    run_all(args.input, args.streams, args.exe)

